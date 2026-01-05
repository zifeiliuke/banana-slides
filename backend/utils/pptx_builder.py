"""
PPTX Builder - utilities for creating editable PPTX files
Based on OpenDCAI/DataFlow-Agent's implementation
"""
import os
import logging
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageFont, ImageDraw
from html.parser import HTMLParser

logger = logging.getLogger(__name__)


class HTMLTableParser(HTMLParser):
    """Parse HTML table into row/column data"""
    
    def __init__(self):
        super().__init__()
        self.table_data = []
        self.current_row = []
        self.current_cell = []
        self.in_table = False
        self.in_row = False
        self.in_cell = False
    
    def handle_starttag(self, tag, attrs):
        if tag == 'table':
            self.in_table = True
            self.table_data = []
        elif tag == 'tr':
            self.in_row = True
            self.current_row = []
        elif tag in ['td', 'th']:
            self.in_cell = True
            self.current_cell = []
    
    def handle_endtag(self, tag):
        if tag == 'table':
            self.in_table = False
        elif tag == 'tr':
            self.in_row = False
            if self.current_row:
                self.table_data.append(self.current_row)
        elif tag in ['td', 'th']:
            self.in_cell = False
            cell_text = ''.join(self.current_cell).strip()
            self.current_row.append(cell_text)
    
    def handle_data(self, data):
        if self.in_cell:
            self.current_cell.append(data)
    
    @staticmethod
    def parse_html_table(html: str) -> List[List[str]]:
        """Parse HTML table string into 2D array of cells"""
        parser = HTMLTableParser()
        parser.feed(html)
        return parser.table_data


class PPTXBuilder:
    """Builder class for creating editable PPTX files from structured content"""
    
    # Standard slide dimensions (16:9 aspect ratio)
    DEFAULT_SLIDE_WIDTH_INCHES = 10
    DEFAULT_SLIDE_HEIGHT_INCHES = 5.625
    
    # Default DPI for pixel to inch conversion
    DEFAULT_DPI = 96
    
    # python-pptx size limits (1-56 inches, 914400-51206400 EMU)
    # See: https://github.com/scanny/python-pptx/issues/93
    MAX_SLIDE_WIDTH_INCHES = 56.0
    MAX_SLIDE_HEIGHT_INCHES = 56.0
    MIN_SLIDE_WIDTH_INCHES = 1.0
    MIN_SLIDE_HEIGHT_INCHES = 1.0
    
    # Global font size limits (to prevent extreme cases)
    MIN_FONT_SIZE = 6   # Minimum readable size
    MAX_FONT_SIZE = 200  # Maximum reasonable size
    
    # 项目内置字体（Noto Sans CJK SC，支持中日韩文字）
    FONT_PATH = os.path.join(os.path.dirname(__file__), "..", "fonts", "NotoSansSC-Regular.ttf")
    
    # Font cache: {size_pt: ImageFont}
    _font_cache: Dict[float, ImageFont.FreeTypeFont] = {}
    
    @classmethod
    def _get_font(cls, size_pt: float) -> Optional[ImageFont.FreeTypeFont]:
        """Get font object for given size (with caching)"""
        # Round to 0.5pt for cache efficiency
        cache_key = round(size_pt * 2) / 2
        
        if cache_key not in cls._font_cache:
            try:
                cls._font_cache[cache_key] = ImageFont.truetype(cls.FONT_PATH, int(size_pt))
            except Exception as e:
                logger.warning(f"Failed to load font {cls.FONT_PATH}: {e}")
                return None
        
        return cls._font_cache[cache_key]
    
    @classmethod
    def _measure_text_width(cls, text: str, font_size_pt: float) -> Optional[float]:
        """
        Measure text width in points using the actual font
        
        Args:
            text: Text to measure
            font_size_pt: Font size in points
            
        Returns:
            Text width in points, or None if measurement failed
        """
        font = cls._get_font(font_size_pt)
        if font is None:
            return None
        
        try:
            # Get text bounding box: (left, top, right, bottom)
            bbox = font.getbbox(text)
            width_px = bbox[2] - bbox[0]
            # Font is loaded at size=font_size_pt, so pixel width ≈ point width
            return width_px
        except Exception as e:
            logger.warning(f"Failed to measure text: {e}")
            return None
    
    def __init__(self, slide_width_inches: float = None, slide_height_inches: float = None):
        """
        Initialize PPTX builder
        
        Args:
            slide_width_inches: Slide width in inches (default: 10)
            slide_height_inches: Slide height in inches (default: 5.625)
        """
        self.slide_width_inches = slide_width_inches or self.DEFAULT_SLIDE_WIDTH_INCHES
        self.slide_height_inches = slide_height_inches or self.DEFAULT_SLIDE_HEIGHT_INCHES
        self.prs = None
        self.current_slide = None
        
    def create_presentation(self) -> Presentation:
        """Create a new presentation with configured dimensions"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.slide_width_inches)
        self.prs.slide_height = Inches(self.slide_height_inches)
        return self.prs
    
    def setup_presentation_size(self, width_pixels: int, height_pixels: int, dpi: int = None):
        """
        Setup presentation size based on pixel dimensions
        Automatically clamps to python-pptx limits (1-56 inches) while preserving aspect ratio
        
        Args:
            width_pixels: Width in pixels
            height_pixels: Height in pixels
            dpi: DPI for conversion (default: 96)
        """
        dpi = dpi or self.DEFAULT_DPI
        
        # Convert pixels to inches
        width_inches = width_pixels / dpi
        height_inches = height_pixels / dpi
        
        # Check if dimensions exceed python-pptx limits and scale down if needed
        # python-pptx enforces: 1 <= dimension <= 56 inches
        scale_factor = 1.0
        
        if width_inches > self.MAX_SLIDE_WIDTH_INCHES:
            scale_factor = self.MAX_SLIDE_WIDTH_INCHES / width_inches
            logger.warning(
                f"Slide width {width_inches:.2f}\" exceeds python-pptx limit ({self.MAX_SLIDE_WIDTH_INCHES}\"), "
                f"scaling down by {scale_factor:.3f}x to maintain aspect ratio"
            )
        
        if height_inches > self.MAX_SLIDE_HEIGHT_INCHES:
            height_scale = self.MAX_SLIDE_HEIGHT_INCHES / height_inches
            if height_scale < scale_factor:
                scale_factor = height_scale
                logger.warning(
                    f"Slide height {height_inches:.2f}\" exceeds python-pptx limit ({self.MAX_SLIDE_HEIGHT_INCHES}\"), "
                    f"scaling down by {scale_factor:.3f}x to maintain aspect ratio"
                )
        
        # Apply scale factor if needed
        if scale_factor < 1.0:
            width_inches *= scale_factor
            height_inches *= scale_factor
            logger.info(
                f"Final slide dimensions after scaling: {width_inches:.2f}\" x {height_inches:.2f}\" "
                f"(from {width_pixels}x{height_pixels}px @ {dpi} DPI)"
            )
        
        # Ensure minimum size constraints
        width_inches = max(self.MIN_SLIDE_WIDTH_INCHES, width_inches)
        height_inches = max(self.MIN_SLIDE_HEIGHT_INCHES, height_inches)
        
        self.slide_width_inches = width_inches
        self.slide_height_inches = height_inches
        
        if self.prs:
            self.prs.slide_width = Inches(self.slide_width_inches)
            self.prs.slide_height = Inches(self.slide_height_inches)
    
    def add_blank_slide(self):
        """Add a blank slide to the presentation"""
        if not self.prs:
            self.create_presentation()
        
        # Use blank layout (layout 6 is typically blank)
        blank_layout = self.prs.slide_layouts[6]
        self.current_slide = self.prs.slides.add_slide(blank_layout)
        return self.current_slide
    
    def pixels_to_inches(self, pixels: float, dpi: int = None) -> float:
        """
        Convert pixels to inches
        
        Args:
            pixels: Pixel value
            dpi: DPI for conversion (default: 96)
            
        Returns:
            Value in inches
        """
        dpi = dpi or self.DEFAULT_DPI
        return pixels / dpi
    
    def calculate_font_size(self, bbox: List[int], text: str, text_level: Any = None, dpi: int = None) -> float:
        """
        Calculate appropriate font size based on bounding box and text content.
        Uses precise font measurement when available, falls back to estimation otherwise.
        Supports both single-line and multi-line (auto-wrap) text.
        
        Args:
            bbox: Bounding box [x0, y0, x1, y1] in pixels
            text: Text content
            text_level: Text level (kept for compatibility, not used in calculation)
            dpi: DPI for pixel to inch conversion
            
        Returns:
            Font size in points (float for precision)
        """
        dpi = dpi or self.DEFAULT_DPI
        
        # Get bbox dimensions in pixels
        width_px = bbox[2] - bbox[0]
        height_px = bbox[3] - bbox[1]
        
        # Convert to points (1 inch = 72 points)
        width_pt = (width_px / dpi) * 72
        height_pt = (height_px / dpi) * 72
        
        # MinerU bbox is tight, use it directly
        # Textbox margins are set to 0 in add_text_element()
        usable_width_pt = width_pt
        usable_height_pt = height_pt
        
        if usable_width_pt <= 0 or usable_height_pt <= 0:
            logger.warning(f"Bbox too small for text: {width_px}x{height_px}px, text: '{text[:30]}...'")
            return self.MIN_FONT_SIZE
        
        text_length = len(text)
        
        # Line height ratio: 1.0 for tight bbox
        line_height_ratio = 1.0
        
        # Try precise measurement first (check if font file exists)
        use_precise = os.path.exists(self.FONT_PATH)
        
        # Binary search: find largest font size that fits
        best_size = self.MIN_FONT_SIZE
        
        for font_size in range(int(self.MAX_FONT_SIZE), int(self.MIN_FONT_SIZE) - 1, -1):
            font_size = float(font_size)
            
            # For text with explicit newlines, calculate each line's width separately
            lines = text.split('\n')
            total_required_lines = 0
            
            for line in lines:
                if not line:
                    total_required_lines += 1
                    continue
                    
                # Measure line width (precise or estimated)
                if use_precise:
                    line_width_pt = self._measure_text_width(line, font_size)
                    if line_width_pt is None:
                        use_precise = False
                
                if not use_precise:
                    # Fallback: estimate based on character count
                    cjk_count = sum(1 for c in line if '\u4e00' <= c <= '\u9fff' or '\u3040' <= c <= '\u30ff' or '\uac00' <= c <= '\ud7af')
                    non_cjk_count = len(line) - cjk_count
                    line_width_pt = (cjk_count * 1.0 + non_cjk_count * 0.5) * font_size
                
                # How many lines does this explicit line need (auto-wrap)?
                lines_needed = max(1, -(-int(line_width_pt) // int(usable_width_pt)))
                total_required_lines += lines_needed
            
            required_lines = total_required_lines
            
            # Calculate total height needed
            line_height_pt = font_size * line_height_ratio
            total_height_pt = required_lines * line_height_pt
            
            # Check if it fits
            if total_height_pt <= usable_height_pt:
                best_size = font_size
                break
        
        if best_size == self.MIN_FONT_SIZE and text_length > 3:
            logger.warning(f"Text may overflow: '{text[:50]}...' in bbox {width_px}x{height_px}px")
        
        # Debug log for font size calculation
        logger.debug(
            f"Font size calc: '{text[:20]}{'...' if len(text) > 20 else ''}' "
            f"bbox={width_px}x{height_px}px -> {best_size}pt "
            f"(usable: {usable_width_pt:.1f}x{usable_height_pt:.1f}pt)"
        )
        
        return best_size
    
    def add_text_element(
        self,
        slide,
        text: str,
        bbox: List[int],
        text_level: Any = None,
        dpi: int = None,
        align: str = 'left',
        text_style: Any = None
    ):
        """
        Add text element to slide
        
        Args:
            slide: Target slide
            text: Text content (used as fallback if text_style has no colored_segments)
            bbox: Bounding box [x0, y0, x1, y1] in pixels
            text_level: Text level (1=title, 2=heading, etc.) or type string
            dpi: DPI for conversion (default: 96)
            align: Text alignment ('left', 'center', 'right')
            text_style: TextStyleResult object with font color, bold, italic etc. (optional)
                        If text_style has colored_segments, those will be used for rendering
                        and the text content will come from the segments.
        """
        dpi = dpi or self.DEFAULT_DPI
        
        # Check if we have colored segments (multi-color text)
        has_colored_segments = (
            text_style and 
            hasattr(text_style, 'colored_segments') and 
            text_style.colored_segments and 
            len(text_style.colored_segments) > 0
        )
        
        # Determine the actual text to use
        # If we have colored_segments, use the text from segments (model's recognized text)
        if has_colored_segments:
            actual_text = ''.join(seg.text for seg in text_style.colored_segments)
        else:
            actual_text = text
        
        # Expand bbox slightly to prevent text overflow
        # MinerU bbox is tight, but font rendering may need extra space
        EXPAND_RATIO = 0.01  # 1% expansion
        bbox_width = bbox[2] - bbox[0]
        bbox_height = bbox[3] - bbox[1]
        expand_w = bbox_width * EXPAND_RATIO
        expand_h = bbox_height * EXPAND_RATIO
        
        # Convert expanded bbox to inches (expand evenly on all sides)
        left = Inches(self.pixels_to_inches(bbox[0] - expand_w / 2, dpi))
        top = Inches(self.pixels_to_inches(bbox[1] - expand_h / 2, dpi))
        width = Inches(self.pixels_to_inches(bbox_width + expand_w, dpi))
        height = Inches(self.pixels_to_inches(bbox_height + expand_h, dpi))
        
        # Add text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # Remove margins completely - bbox is tight, no extra space needed
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0)
        text_frame.margin_top = Inches(0)
        text_frame.margin_bottom = Inches(0)
        
        def replace_some_chars(text: str) -> str:
            # replace logic
            # replace · to • if starts with ·
            text = text.replace('·', '•', 1) if text.lstrip().startswith('·') else text
            return text
        actual_text = replace_some_chars(actual_text)
        
        # Calculate font size
        font_size = self.calculate_font_size(bbox, actual_text, text_level, dpi)
        
        # Determine effective alignment - text_style优先，否则使用参数
        effective_align = align
        if text_style and hasattr(text_style, 'text_alignment') and text_style.text_alignment:
            effective_align = text_style.text_alignment
        
        # Get style attributes
        is_bold = False
        is_italic = False
        is_underline = False
        if text_style:
            is_bold = getattr(text_style, 'is_bold', False)
            is_italic = getattr(text_style, 'is_italic', False)
            is_underline = getattr(text_style, 'is_underline', False)
        
        # Make title text bold (legacy behavior)
        if text_level == 1 or text_level == 'title':
            is_bold = True
        
        # Render text with colors
        if has_colored_segments:
            # Multi-color text: use runs for each segment
            paragraph = text_frame.paragraphs[0]
            paragraph.clear()
            
            latex_count = 0
            for seg in text_style.colored_segments:
                run = paragraph.add_run()
                run.text = replace_some_chars(seg.text)
                run.font.size = Pt(font_size)
                run.font.bold = is_bold
                run.font.underline = is_underline
                # Set segment-specific color
                r, g, b = seg.color_rgb
                run.font.color.rgb = RGBColor(r, g, b)
                
                # Handle LaTeX formula segments
                if hasattr(seg, 'is_latex') and seg.is_latex:
                    # For LaTeX formulas, use italic style as visual hint
                    # TODO: In future, could render as actual equation using OMML
                    run.font.italic = True
                    latex_count += 1
                    logger.debug(f"  LaTeX formula detected: '{seg.text}'")
                else:
                    run.font.italic = is_italic
            
            latex_info = f", {latex_count} latex" if latex_count > 0 else ""
            style_info = f" | multi-color: {len(text_style.colored_segments)} segments{latex_info}"
        else:
            # Single color text: use simple text assignment
            text_frame.text = actual_text
            # IMPORTANT: Re-get paragraph after setting text_frame.text
            # because setting text_frame.text creates a new paragraph object
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = is_bold
            paragraph.font.italic = is_italic
            paragraph.font.underline = is_underline
            
            # Apply single font color if provided
            if text_style and hasattr(text_style, 'font_color_rgb') and text_style.font_color_rgb:
                r, g, b = text_style.font_color_rgb
                paragraph.font.color.rgb = RGBColor(r, g, b)
            
            style_info = f" | color={text_style.font_color_rgb if text_style else 'default'}"
        
        # Apply alignment after paragraph is finalized
        if effective_align == 'center':
            paragraph.alignment = PP_ALIGN.CENTER
        elif effective_align == 'right':
            paragraph.alignment = PP_ALIGN.RIGHT
        elif effective_align == 'justify':
            paragraph.alignment = PP_ALIGN.JUSTIFY
        else:
            paragraph.alignment = PP_ALIGN.LEFT
        
        # Calculate bbox dimensions for logging
        bbox_width = bbox[2] - bbox[0]
        bbox_height = bbox[3] - bbox[1]
        logger.debug(f"Text: '{actual_text[:35]}' | box: {bbox_width}x{bbox_height}px | font: {font_size:.1f}pt | chars: {len(actual_text)}{style_info}")
    
    def add_image_element(
        self,
        slide,
        image_path: str,
        bbox: List[int],
        dpi: int = None
    ):
        """
        Add image element to slide
        
        Args:
            slide: Target slide
            image_path: Path to image file
            bbox: Bounding box [x0, y0, x1, y1] in pixels
            dpi: DPI for conversion (default: 96)
        """
        dpi = dpi or self.DEFAULT_DPI
        
        # Check if image exists
        if not os.path.exists(image_path):
            logger.warning(f"Image not found: {image_path}, adding placeholder")
            self.add_image_placeholder(slide, bbox, dpi)
            return
        
        # Convert bbox to inches
        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))
        
        try:
            # Add image
            slide.shapes.add_picture(image_path, left, top, width, height)
            logger.debug(f"Added image: {image_path} at bbox {bbox}")
        except Exception as e:
            logger.error(f"Failed to add image {image_path}: {str(e)}")
            self.add_image_placeholder(slide, bbox, dpi)
    
    def add_image_placeholder(
        self,
        slide,
        bbox: List[int],
        dpi: int = None
    ):
        """
        Add a placeholder for missing images
        
        Args:
            slide: Target slide
            bbox: Bounding box [x0, y0, x1, y1] in pixels
            dpi: DPI for conversion (default: 96)
        """
        dpi = dpi or self.DEFAULT_DPI
        
        # Convert bbox to inches
        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))
        
        # Add a text box as placeholder
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "[Image]"
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(12)
        paragraph.font.italic = True
    
    def add_table_element(
        self,
        slide,
        html_table: str,
        bbox: List[int],
        dpi: int = None
    ):
        """
        Add editable table to slide from HTML table string
        
        Args:
            slide: Target slide
            html_table: HTML table string
            bbox: Bounding box [x0, y0, x1, y1] in pixels
            dpi: DPI for conversion (default: 96)
        """
        dpi = dpi or self.DEFAULT_DPI
        
        # Parse HTML table
        try:
            table_data = HTMLTableParser.parse_html_table(html_table)
        except Exception as e:
            logger.error(f"Failed to parse HTML table: {str(e)}")
            return
        
        if not table_data or not table_data[0]:
            logger.warning("Empty table data")
            return
        
        rows = len(table_data)
        cols = len(table_data[0])
        
        # Convert bbox to inches
        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))
        
        try:
            # Add table shape
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            
            # Calculate cell dimensions
            cell_width = width / cols
            cell_height = height / rows
            
            # Fill table with data
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < cols:  # Safety check
                        cell = table.cell(row_idx, col_idx)
                        cell.text = cell_text
                        
                        # Style the cell
                        text_frame = cell.text_frame
                        text_frame.word_wrap = True
                        
                        # Calculate font size for table cell
                        # Use a conservative size to fit in cell
                        cell_height_px = (bbox[3] - bbox[1]) / rows
                        cell_width_px = (bbox[2] - bbox[0]) / cols
                        
                        # Estimate font size (smaller for tables)
                        font_size = min(18, max(8, cell_height_px * 0.3))
                        
                        for paragraph in text_frame.paragraphs:
                            paragraph.font.size = Pt(font_size)
                            paragraph.alignment = PP_ALIGN.CENTER
                            
                            # Header row (first row) should be bold
                            if row_idx == 0:
                                paragraph.font.bold = True
            
            logger.info(f"Added editable table: {rows}x{cols} at bbox {bbox}")
            
        except Exception as e:
            logger.error(f"Failed to create table: {str(e)}")
    
    def save(self, output_path: str):
        """
        Save presentation to file
        
        Args:
            output_path: Output file path
        """
        if not self.prs:
            raise ValueError("No presentation to save")
        
        # Ensure directory exists
        output_path_obj = Path(output_path)
        output_dir = output_path_obj.parent
        if str(output_dir) != '.':  # Only create directory if it's not current directory
            output_dir.mkdir(parents=True, exist_ok=True)
        
        self.prs.save(output_path)
        logger.info(f"Saved presentation to: {output_path}")
    
    def get_presentation(self) -> Presentation:
        """Get the current presentation object"""
        return self.prs

