"""
Export Service - handles PPTX and PDF export
Based on demo.py create_pptx_from_images()
"""
import os
import json
import logging
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from textwrap import dedent
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io
import tempfile
import img2pdf
logger = logging.getLogger(__name__)


@dataclass
class ExportWarnings:
    """
    导出过程中收集的警告信息
    
    用于追踪哪些操作没有按预期执行，并反馈给前端
    """
    # 样式提取失败的元素
    style_extraction_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # 文本渲染失败的元素
    text_render_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # 图片添加失败
    image_add_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # JSON 解析失败（重试后仍失败）
    json_parse_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # 其他警告
    other_warnings: List[str] = field(default_factory=list)
    
    def add_style_extraction_failed(self, element_id: str, reason: str):
        """记录样式提取失败"""
        self.style_extraction_failed.append({
            'element_id': element_id,
            'reason': reason
        })
    
    def add_text_render_failed(self, text: str, reason: str):
        """记录文本渲染失败"""
        self.text_render_failed.append({
            'text': text[:50] + '...' if len(text) > 50 else text,
            'reason': reason
        })
    
    def add_image_failed(self, path: str, reason: str):
        """记录图片添加失败"""
        self.image_add_failed.append({
            'path': path,
            'reason': reason
        })
    
    def add_json_parse_failed(self, context: str, reason: str):
        """记录 JSON 解析失败"""
        self.json_parse_failed.append({
            'context': context,
            'reason': reason
        })
    
    def add_warning(self, message: str):
        """添加其他警告"""
        self.other_warnings.append(message)
    
    def has_warnings(self) -> bool:
        """是否有警告"""
        return bool(
            self.style_extraction_failed or 
            self.text_render_failed or 
            self.image_add_failed or
            self.json_parse_failed or
            self.other_warnings
        )
    
    def to_summary(self) -> List[str]:
        """生成警告摘要（适合前端展示）"""
        summary = []
        
        if self.style_extraction_failed:
            summary.append(f"⚠️ {len(self.style_extraction_failed)} 个文本元素样式提取失败")
        
        if self.text_render_failed:
            summary.append(f"⚠️ {len(self.text_render_failed)} 个文本元素渲染失败")
        
        if self.image_add_failed:
            summary.append(f"⚠️ {len(self.image_add_failed)} 张图片添加失败")
        
        if self.json_parse_failed:
            summary.append(f"⚠️ {len(self.json_parse_failed)} 次 AI 响应解析失败")
        
        for warning in self.other_warnings[:5]:  # 最多显示5条其他警告
            summary.append(f"⚠️ {warning}")
        
        if len(self.other_warnings) > 5:
            summary.append(f"  ...还有 {len(self.other_warnings) - 5} 条其他警告")
        
        return summary
    
    def to_dict(self) -> Dict[str, Any]:
        """转换为字典（详细信息）"""
        return {
            'style_extraction_failed': self.style_extraction_failed,
            'text_render_failed': self.text_render_failed,
            'image_add_failed': self.image_add_failed,
            'json_parse_failed': self.json_parse_failed,
            'other_warnings': self.other_warnings,
            'total_warnings': (
                len(self.style_extraction_failed) + 
                len(self.text_render_failed) + 
                len(self.image_add_failed) +
                len(self.json_parse_failed) +
                len(self.other_warnings)
            )
        }


class ExportService:
    """Service for exporting presentations"""
    
    # NOTE: clean background生成功能已迁移到解耦的InpaintProvider实现
    # - DefaultInpaintProvider: 基于mask的精确区域重绘（Volcengine）
    # - GenerativeEditInpaintProvider: 基于生成式大模型的整图编辑重绘（Gemini等）
    # 使用方式: from services.image_editability import InpaintProviderFactory
    
    @staticmethod
    def create_pptx_from_images(image_paths: List[str], output_file: str = None) -> bytes:
        """
        Create PPTX file from image paths
        Based on demo.py create_pptx_from_images()
        
        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)
        
        Returns:
            PPTX file as bytes if output_file is None
        """
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions to 16:9 (width 10 inches, height 5.625 inches)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # Add each image as a slide
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue
            
            # Add blank slide layout (layout 6 is typically blank)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to fill entire slide
            slide.shapes.add_picture(
                image_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
        
        # Save or return bytes
        if output_file:
            prs.save(output_file)
            return None
        else:
            # Save to bytes
            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes.getvalue()
    
    @staticmethod
    def create_pdf_from_images(image_paths: List[str], output_file: str = None) -> Optional[bytes]:
        """
        Create PDF file from image paths using img2pdf (low memory usage)

        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)

        Returns:
            PDF file as bytes if output_file is None, otherwise None
        """
        # Validate images exist and log warnings for missing files
        valid_paths = []
        for p in image_paths:
            if os.path.exists(p):
                valid_paths.append(p)
            else:
                logger.warning(f"Image not found and will be skipped for PDF export: {p}")

        if not valid_paths:
            raise ValueError("No valid images found for PDF export")

        try:
            logger.info(f"Using img2pdf for PDF export ({len(valid_paths)} pages, low memory mode)")

            # Set page layout: 16:9 aspect ratio (10 inches × 5.625 inches)
            layout_fun = img2pdf.get_layout_fun(
                pagesize=(img2pdf.in_to_pt(10), img2pdf.in_to_pt(5.625))
            )

            # Convert images to PDF
            pdf_bytes = img2pdf.convert(valid_paths, layout_fun=layout_fun)

            if output_file:
                with open(output_file, "wb") as f:
                    f.write(pdf_bytes)
                return None
            else:
                return pdf_bytes
        except (img2pdf.ImageOpenError, ValueError, IOError) as e:
            logger.warning(f"img2pdf conversion failed: {e}. Falling back to Pillow (high memory usage).")
            return ExportService.create_pdf_from_images_pillow(valid_paths, output_file)

    @staticmethod
    def create_pdf_from_images_pillow(image_paths: List[str], output_file: str = None) -> Optional[bytes]:
        """
        Create PDF file from image paths using Pillow (original method)

        Note: This method loads all images into memory at once.
        For large projects (50+ pages with 20MB/page), use create_pdf_from_images instead.

        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)

        Returns:
            PDF file as bytes if output_file is None, otherwise None
        """
        images = []

        # Load all images
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue

            img = Image.open(image_path)

            # Convert to RGB if necessary (PDF requires RGB)
            if img.mode != 'RGB':
                img = img.convert('RGB')

            images.append(img)

        if not images:
            raise ValueError("No valid images found for PDF export")

        # Save as PDF
        if output_file:
            images[0].save(
                output_file,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            return None
        else:
            # Save to bytes
            pdf_bytes = io.BytesIO()
            images[0].save(
                pdf_bytes,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            pdf_bytes.seek(0)
            return pdf_bytes.getvalue()
       
    @staticmethod
    def _add_mineru_text_to_slide(builder, slide, text_item: Dict[str, Any], scale_x: float = 1.0, scale_y: float = 1.0):
        """
        Add text item from MinerU to slide
        
        Args:
            builder: PPTXBuilder instance
            slide: Target slide
            text_item: Text item from MinerU content_list
            scale_x: X-axis scale factor
            scale_y: Y-axis scale factor
        """
        text = text_item.get('text', '').strip()
        if not text:
            return
        
        bbox = text_item.get('bbox')
        if not bbox or len(bbox) != 4:
            logger.warning(f"Invalid bbox for text item: {text_item}")
            return
        
        original_bbox = bbox.copy()
        
        # Apply scale factors to bbox
        x0, y0, x1, y1 = bbox
        bbox = [
            int(x0 * scale_x),
            int(y0 * scale_y),
            int(x1 * scale_x),
            int(y1 * scale_y)
        ]
        
        if scale_x != 1.0 or scale_y != 1.0:
            logger.debug(f"Text bbox scaled: {original_bbox} -> {bbox} (scale: {scale_x:.3f}x{scale_y:.3f})")
        
        # Determine text level (only used for styling like bold, NOT for font size)
        # Font size is purely calculated from bbox dimensions
        item_type = text_item.get('type', 'text')
        text_level = text_item.get('text_level')
        
        # Map to level for styling purposes (bold titles)
        if item_type == 'title' or text_level == 1:
            level = 'title'  # Will be bold
        else:
            level = 'default'
        
        # Add text element
        # Note: text_level is only used for bold styling, not font size calculation
        try:
            builder.add_text_element(
                slide=slide,
                text=text,
                bbox=bbox,
                text_level=level  # For styling (bold) only, not font size
            )
        except Exception as e:
            logger.error(f"Failed to add text element: {str(e)}")
    
    @staticmethod
    def _add_table_cell_elements_to_slide(
        builder,
        slide,
        cell_elements: List[Dict[str, Any]],
        scale_x: float = 1.0,
        scale_y: float = 1.0
    ):
        """
        Add table cell elements as individual text boxes to slide
        这些单元格元素已经有正确的全局bbox坐标
        
        Args:
            builder: PPTXBuilder instance
            slide: Target slide
            cell_elements: List of EditableElement (table_cell type)
            scale_x: X-axis scale factor
            scale_y: Y-axis scale factor
        """
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        
        logger.info(f"开始添加表格单元格元素，共 {len(cell_elements)} 个")
        
        for cell_elem in cell_elements:
            text = cell_elem.get('content', '')
            bbox_global = cell_elem.get('bbox_global', {})
            
            if not text.strip():
                continue
            
            # bbox_global已经是全局坐标，直接使用并应用缩放
            x0 = bbox_global.get('x0', 0)
            y0 = bbox_global.get('y0', 0)
            x1 = bbox_global.get('x1', 0)
            y1 = bbox_global.get('y1', 0)
            
            # 构建bbox列表 [x0, y0, x1, y1] 并应用缩放
            bbox = [
                int(x0 * scale_x),
                int(y0 * scale_y),
                int(x1 * scale_x),
                int(y1 * scale_y)
            ]
            
            try:
                # 使用已有的 add_text_element 方法添加文本框（不添加边框）
                builder.add_text_element(
                    slide=slide,
                    text=text,
                    bbox=bbox,
                    text_level=None,
                    align='center'
                )
                
                logger.debug(f"  添加单元格: '{text[:10]}...' at bbox {bbox}")
                
            except Exception as e:
                logger.warning(f"添加单元格失败: {e}")
        
        logger.info(f"✓ 表格单元格添加完成，共 {len(cell_elements)} 个")
    
    @staticmethod
    def _add_mineru_image_to_slide(
        builder,
        slide,
        image_item: Dict[str, Any],
        mineru_dir: Path,
        scale_x: float = 1.0,
        scale_y: float = 1.0
    ):
        """
        Add image or table item from MinerU to slide
        
        Args:
            builder: PPTXBuilder instance
            slide: Target slide
            image_item: Image/table item from MinerU content_list
            mineru_dir: MinerU result directory
            scale_x: X-axis scale factor
            scale_y: Y-axis scale factor
        """
        bbox = image_item.get('bbox')
        if not bbox or len(bbox) != 4:
            logger.warning(f"Invalid bbox for image item: {image_item}")
            return
        
        original_bbox = bbox.copy()
        
        # Apply scale factors to bbox
        x0, y0, x1, y1 = bbox
        bbox = [
            int(x0 * scale_x),
            int(y0 * scale_y),
            int(x1 * scale_x),
            int(y1 * scale_y)
        ]
        
        if scale_x != 1.0 or scale_y != 1.0:
            logger.debug(f"Item bbox scaled: {original_bbox} -> {bbox} (scale: {scale_x:.3f}x{scale_y:.3f})")
        
        # Check if this is a table with子元素 (cells from Baidu OCR)
        item_type = image_item.get('element_type') or image_item.get('type', 'image')
        children = image_item.get('children', [])
        
        logger.debug(f"Processing {item_type} element, has {len(children)} children")
        
        if children and item_type == 'table':
            # Add editable table from child elements (cells)
            try:
                # Filter only table_cell elements
                cell_elements = [child for child in children if child.get('element_type') == 'table_cell']
                
                if cell_elements:
                    logger.info(f"添加可编辑表格（{len(cell_elements)}个单元格）")
                    ExportService._add_table_cell_elements_to_slide(
                        builder=builder,
                        slide=slide,
                        cell_elements=cell_elements,
                        scale_x=scale_x,
                        scale_y=scale_y
                    )
                    return  # Table added successfully
            except Exception as e:
                logger.exception("Failed to add table cells, falling back to image")
                # Fall through to add as image instead
        
        # Check if this is a table with HTML data (legacy)
        html_table = image_item.get('html_table')
        if html_table and item_type == 'table':
            # Add editable table from HTML
            try:
                builder.add_table_element(
                    slide=slide,
                    html_table=html_table,
                    bbox=bbox
                )
                logger.info(f"Added editable table at bbox {bbox}")
                return  # Table added successfully
            except Exception as e:
                logger.error(f"Failed to add table: {str(e)}, falling back to image")
                # Fall through to add as image instead
        
        # Add as image (either image type or table fallback)
        img_path_str = image_item.get('img_path', '')
        if not img_path_str:
            logger.warning(f"No img_path in item: {image_item}")
            return
        
        # Try to find the image file
        # MinerU may store images in 'images/' subdirectory
        possible_paths = [
            mineru_dir / img_path_str,
            mineru_dir / 'images' / Path(img_path_str).name,
            mineru_dir / Path(img_path_str).name,
        ]
        
        image_path = None
        for path in possible_paths:
            if path.exists():
                image_path = str(path)
                break
        
        if not image_path:
            logger.warning(f"Image file not found: {img_path_str}")
            # Add placeholder
            builder.add_image_placeholder(slide, bbox)
            return
        
        # Add image element
        try:
            builder.add_image_element(
                slide=slide,
                image_path=image_path,
                bbox=bbox
            )
        except Exception as e:
            logger.error(f"Failed to add image element: {str(e)}")
    
    @staticmethod
    def _collect_text_elements_for_extraction(
        elements: List,  # List[EditableElement]
        depth: int = 0
    ) -> List[tuple]:
        """
        递归收集所有需要提取样式的文本元素
        
        Args:
            elements: EditableElement列表
            depth: 当前递归深度
        
        Returns:
            元组列表，每个元组为 (element_id, image_path, text_content)
        """
        text_items = []
        
        for elem in elements:
            elem_type = elem.element_type
            
            # 文本类型元素需要提取样式
            if elem_type in ['text', 'title', 'table_cell', 'list', 'paragraph', 'header', 'footer', 'heading', 'table_caption', 'image_caption']:
                if elem.content and elem.image_path and os.path.exists(elem.image_path):
                    text = elem.content.strip()
                    if text:
                        text_items.append((elem.element_id, elem.image_path, text))
            
            # 递归处理子元素
            if hasattr(elem, 'children') and elem.children:
                child_items = ExportService._collect_text_elements_for_extraction(
                    elements=elem.children,
                    depth=depth + 1
                )
                text_items.extend(child_items)
        
        return text_items
    
    @staticmethod
    def _batch_extract_text_styles(
        text_items: List[tuple],
        text_attribute_extractor,
        max_workers: int = 8
    ) -> Dict[str, Any]:
        """
        批量并行提取文本样式（逐个裁剪区域分析）
        
        此方法对每一段文字的裁剪区域单独进行分析。
        经测试，此方法效果较好，目前仍在使用。
        
        备选方案：_batch_extract_text_styles_with_full_image 可一次性分析全图所有文本。
        
        Args:
            text_items: 元组列表，每个元组为 (element_id, image_path, text_content)
            text_attribute_extractor: 文本属性提取器
            max_workers: 并发数
        
        Returns:
            字典，key为element_id，value为TextStyleResult
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        
        if not text_items or not text_attribute_extractor:
            return {}
        
        logger.info(f"并行提取 {len(text_items)} 个文本元素的样式（并发数: {max_workers}）...")
        
        results = {}
        
        def extract_single(item):
            element_id, image_path, text_content = item
            try:
                style = text_attribute_extractor.extract(
                    image=image_path,
                    text_content=text_content
                )
                return element_id, style
            except Exception as e:
                logger.warning(f"提取文字样式失败 [{element_id}]: {e}")
                return element_id, None
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(extract_single, item): item[0] for item in text_items}
            
            for future in as_completed(futures):
                element_id, style = future.result()
                if style is not None:
                    results[element_id] = style
        
        logger.info(f"✓ 文本样式提取完成，成功 {len(results)}/{len(text_items)} 个")
        return results
    
    @staticmethod
    def _collect_text_elements_for_batch_extraction(
        elements: List,  # List[EditableElement]
        depth: int = 0
    ) -> List[Dict[str, Any]]:
        """
        递归收集所有需要批量提取样式的文本元素（新格式，包含bbox）
        
        Args:
            elements: EditableElement列表
            depth: 当前递归深度
        
        Returns:
            字典列表，每个字典包含 element_id, bbox, content
        """
        text_items = []
        
        for elem in elements:
            elem_type = elem.element_type
            
            # 文本类型元素需要提取样式
            if elem_type in ['text', 'title', 'table_cell', 'list', 'paragraph', 'header', 'footer', 'heading', 'table_caption', 'image_caption']:
                if elem.content:
                    text = elem.content.strip()
                    if text:
                        # 使用全局坐标 bbox_global
                        bbox = elem.bbox_global if hasattr(elem, 'bbox_global') and elem.bbox_global else elem.bbox
                        text_items.append({
                            'element_id': elem.element_id,
                            'bbox': [bbox.x0, bbox.y0, bbox.x1, bbox.y1],
                            'content': text
                        })
            
            # 递归处理子元素
            if hasattr(elem, 'children') and elem.children:
                child_items = ExportService._collect_text_elements_for_batch_extraction(
                    elements=elem.children,
                    depth=depth + 1
                )
                text_items.extend(child_items)
        
        return text_items
    
    @staticmethod
    def _batch_extract_text_styles_with_full_image(
        editable_images: List,  # List[EditableImage]
        text_attribute_extractor,
        max_workers: int = 4
    ) -> Dict[str, Any]:
        """
        【新逻辑】使用全图批量提取所有文本样式
        
        新方法：给 caption model 提供全图，以及提取后的所有文本 bbox 和内容，
        让模型一次性分析所有文本的样式属性（颜色、粗体、对齐等）。
        
        优势：模型可以看到全局信息，分析更准确。
        
        Args:
            editable_images: EditableImage列表，每个对应一张PPT页面
            text_attribute_extractor: 文本属性提取器（需要有 extract_batch_with_full_image 方法）
            max_workers: 并发处理页面数
        
        Returns:
            字典，key为element_id，value为TextStyleResult
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        
        if not editable_images or not text_attribute_extractor:
            return {}
        
        # 检查提取器是否支持批量提取
        if not hasattr(text_attribute_extractor, 'extract_batch_with_full_image'):
            logger.warning("提取器不支持 extract_batch_with_full_image 方法，回退到旧逻辑")
            # 回退到旧逻辑
            all_text_items = []
            for editable_img in editable_images:
                text_items = ExportService._collect_text_elements_for_extraction(editable_img.elements)
                all_text_items.extend(text_items)
            return ExportService._batch_extract_text_styles(
                text_items=all_text_items,
                text_attribute_extractor=text_attribute_extractor,
                max_workers=max_workers * 2
            )
        
        logger.info(f"【新逻辑】使用全图批量分析 {len(editable_images)} 页的文本样式...")
        
        all_results = {}
        
        def process_single_page(editable_img, page_idx):
            """处理单个页面的文本样式提取"""
            try:
                # 收集该页面的所有文本元素
                text_elements = ExportService._collect_text_elements_for_batch_extraction(
                    editable_img.elements
                )
                
                if not text_elements:
                    logger.info(f"  页面 {page_idx + 1}: 无文本元素")
                    return {}
                
                logger.info(f"  页面 {page_idx + 1}: 分析 {len(text_elements)} 个文本元素...")
                
                # 使用原始图片路径作为全图
                full_image_path = editable_img.image_path
                
                # 调用批量提取方法
                page_results = text_attribute_extractor.extract_batch_with_full_image(
                    full_image=full_image_path,
                    text_elements=text_elements
                )
                
                logger.info(f"  页面 {page_idx + 1}: 成功提取 {len(page_results)} 个元素的样式")
                return page_results
                
            except Exception as e:
                logger.error(f"页面 {page_idx + 1} 文本样式提取失败: {e}", exc_info=True)
                return {}
        
        # 并发处理所有页面
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {
                executor.submit(process_single_page, img, idx): idx 
                for idx, img in enumerate(editable_images)
            }
            
            for future in as_completed(futures):
                page_idx = futures[future]
                try:
                    page_results = future.result()
                    all_results.update(page_results)
                except Exception as e:
                    logger.error(f"页面 {page_idx + 1} 处理失败: {e}")
        
        total_elements = sum(
            len(ExportService._collect_text_elements_for_batch_extraction(img.elements))
            for img in editable_images
        )
        logger.info(f"✓ 全图批量文本样式提取完成，成功 {len(all_results)}/{total_elements} 个")
        
        return all_results
    
    @staticmethod
    def _batch_extract_text_styles_hybrid(
        editable_images: List,  # List[EditableImage]
        text_attribute_extractor,
        max_workers: int = 8
    ) -> Tuple[Dict[str, Any], List[Tuple[str, str]]]:
        """
        【混合策略】结合全局识别和单个裁剪识别的优势
        
        策略：
        - 全局识别（全图分析）：获取 is_bold、is_italic、is_underline、text_alignment
          因为这些属性需要看整体布局和上下文才能判断准确
        - 单个裁剪识别：获取 font_color
          因为颜色需要精确看局部像素才能识别准确
        
        Args:
            editable_images: EditableImage列表，每个对应一张PPT页面
            text_attribute_extractor: 文本属性提取器
            max_workers: 并发数
        
        Returns:
            (results, failed_extractions):
            - results: 字典，key为element_id，value为TextStyleResult（合并后的结果）
            - failed_extractions: 失败列表，每项为 (element_id, error_reason)
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        from services.image_editability.text_attribute_extractors import TextStyleResult
        
        if not editable_images or not text_attribute_extractor:
            return {}, []
        
        # 检查提取器是否支持批量提取
        if not hasattr(text_attribute_extractor, 'extract_batch_with_full_image'):
            logger.warning("提取器不支持混合策略，回退到单个裁剪识别")
            all_text_items = []
            for editable_img in editable_images:
                text_items = ExportService._collect_text_elements_for_extraction(editable_img.elements)
                all_text_items.extend(text_items)
            results = ExportService._batch_extract_text_styles(
                text_items=all_text_items,
                text_attribute_extractor=text_attribute_extractor,
                max_workers=max_workers
            )
            return results, []  # 回退方法暂不收集失败信息
        
        logger.info(f"【混合策略】开始分析 {len(editable_images)} 页的文本样式...")
        logger.info(f"  - 全局识别: is_bold, is_italic, is_underline, text_alignment")
        logger.info(f"  - 单个识别: font_color")
        
        # Step 1: 收集所有文本元素
        all_text_items = []  # 用于单个裁剪识别 (element_id, image_path, content)
        page_text_elements = {}  # 用于全局识别 {page_idx: [text_elements]}
        
        for page_idx, editable_img in enumerate(editable_images):
            # 收集用于单个裁剪识别的数据
            text_items = ExportService._collect_text_elements_for_extraction(editable_img.elements)
            all_text_items.extend(text_items)
            
            # 收集用于全局识别的数据
            batch_elements = ExportService._collect_text_elements_for_batch_extraction(editable_img.elements)
            if batch_elements:
                page_text_elements[page_idx] = {
                    'image_path': editable_img.image_path,
                    'elements': batch_elements
                }
        
        if not all_text_items:
            return {}
        
        # Step 2: 并行执行两种识别
        global_results = {}  # 全局识别结果
        local_results = {}   # 单个裁剪识别结果
        
        def extract_global_for_page(page_idx, page_data):
            """全局识别单页"""
            try:
                results = text_attribute_extractor.extract_batch_with_full_image(
                    full_image=page_data['image_path'],
                    text_elements=page_data['elements']
                )
                return page_idx, results
            except Exception as e:
                logger.warning(f"全局识别页面 {page_idx + 1} 失败: {e}")
                return page_idx, {}
        
        # 收集失败信息
        failed_extractions = []  # [(element_id, reason), ...]
        
        def extract_local_single(item):
            """单个裁剪识别"""
            element_id, image_path, text_content = item
            try:
                style = text_attribute_extractor.extract(
                    image=image_path,
                    text_content=text_content
                )
                # 只要 style 不为 None 就算成功（黑色也是有效颜色）
                if style:
                    return element_id, style, None
                else:
                    return element_id, None, "样式提取返回空"
            except Exception as e:
                logger.warning(f"单个识别失败 [{element_id}]: {e}")
                return element_id, None, str(e)
        
        # 并发执行全局识别和单个裁剪识别
        logger.info(f"  并发执行: 全局识别 {len(page_text_elements)} 页 + 单个识别 {len(all_text_items)} 个元素...")
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # 提交全局识别任务
            global_futures = {
                executor.submit(extract_global_for_page, idx, data): ('global', idx)
                for idx, data in page_text_elements.items()
            }
            
            # 提交单个裁剪识别任务
            local_futures = {
                executor.submit(extract_local_single, item): ('local', item[0])
                for item in all_text_items
            }
            
            # 收集全局识别结果
            for future in as_completed(global_futures):
                task_type, page_idx = global_futures[future]
                try:
                    _, page_results = future.result()
                    global_results.update(page_results)
                except Exception as e:
                    logger.error(f"全局识别任务失败: {e}")
            
            # 收集单个裁剪识别结果
            for future in as_completed(local_futures):
                task_type, element_id = local_futures[future]
                try:
                    elem_id, style, error = future.result()
                    if style is not None:
                        local_results[elem_id] = style
                    if error:
                        failed_extractions.append((elem_id, error))
                except Exception as e:
                    logger.error(f"单个识别任务失败: {e}")
                    failed_extractions.append((element_id, str(e)))
        
        # Step 3: 合并结果
        # 优先使用全局识别的布局属性，使用单个识别的颜色属性
        merged_results = {}
        
        all_element_ids = set(global_results.keys()) | set(local_results.keys())
        
        for element_id in all_element_ids:
            global_style = global_results.get(element_id)
            local_style = local_results.get(element_id)
            
            if global_style and local_style:
                # 混合：颜色用单个识别（包括 colored_segments），布局用全局识别
                merged_results[element_id] = TextStyleResult(
                    font_color_rgb=local_style.font_color_rgb,  # 单个识别的颜色
                    colored_segments=local_style.colored_segments,  # 单个识别的多颜色片段
                    is_bold=global_style.is_bold,              # 全局识别的粗体
                    is_italic=global_style.is_italic,          # 全局识别的斜体
                    is_underline=global_style.is_underline,    # 全局识别的下划线
                    text_alignment=global_style.text_alignment, # 全局识别的对齐
                    confidence=0.9,
                    metadata={
                        'source': 'hybrid',
                        'color_source': 'local',
                        'layout_source': 'global'
                    }
                )
            elif local_style:
                # 只有单个识别结果
                merged_results[element_id] = local_style
            elif global_style:
                # 只有全局识别结果
                merged_results[element_id] = global_style
        
        logger.info(f"✓ 混合策略完成: 全局识别 {len(global_results)} 个, 单个识别 {len(local_results)} 个, 合并 {len(merged_results)} 个, 失败 {len(failed_extractions)} 个")
        
        return merged_results, failed_extractions
    
    @staticmethod
    def create_editable_pptx_with_recursive_analysis(
        image_paths: List[str] = None,
        output_file: str = None,
        slide_width_pixels: int = 1920,
        slide_height_pixels: int = 1080,
        max_depth: int = 2,
        max_workers: int = 8,
        editable_images: List = None,  # 可选：直接传入已分析的EditableImage列表
        text_attribute_extractor = None,  # 可选：文字属性提取器，用于提取颜色、粗体、斜体等样式
        progress_callback = None,  # 可选：进度回调函数 (step, message, percent) -> None
        export_extractor_method: str = 'hybrid',  # 组件提取方法: mineru, hybrid
        export_inpaint_method: str = 'hybrid'  # 背景修复方法: generative, baidu, hybrid
    ) -> Tuple[Optional[bytes], ExportWarnings]:
        """
        使用递归图片可编辑化服务创建可编辑PPTX
        
        这是新的架构方法，使用ImageEditabilityService进行递归版面分析。
        
        两种使用方式：
        1. 传入 image_paths：自动分析图片并生成PPTX
        2. 传入 editable_images：直接使用已分析的结果（避免重复分析）
        
        配置（如 MinerU token）自动从 Flask app.config 获取。
        
        Args:
            image_paths: 图片路径列表（可选，与editable_images二选一）
            output_file: 输出文件路径（可选）
            slide_width_pixels: 目标幻灯片宽度
            slide_height_pixels: 目标幻灯片高度
            max_depth: 最大递归深度
            max_workers: 并发处理数
            editable_images: 已分析的EditableImage列表（可选，与image_paths二选一）
            text_attribute_extractor: 文字属性提取器（可选），用于提取文字颜色、粗体、斜体等样式
                可通过 TextAttributeExtractorFactory.create_caption_model_extractor() 创建
            export_extractor_method: 组件提取方法 ('mineru' 或 'hybrid'，默认 'hybrid')
            export_inpaint_method: 背景修复方法 ('generative', 'baidu', 'hybrid'，默认 'hybrid')
        
        Returns:
            (pptx_bytes, warnings): 元组，包含 PPTX 字节流和警告信息
            - pptx_bytes: PPTX 文件字节流（如果 output_file 为 None），否则为 None
            - warnings: ExportWarnings 对象，包含所有警告信息
        """
        from services.image_editability import ServiceConfig, ImageEditabilityService
        from utils.pptx_builder import PPTXBuilder
        
        # 初始化警告收集器
        warnings = ExportWarnings()
        
        # 辅助函数：报告进度
        def report_progress(step: str, message: str, percent: int):
            logger.info(f"[进度 {percent}%] {step}: {message}")
            if progress_callback:
                try:
                    progress_callback(step, message, percent)
                except Exception as e:
                    logger.warning(f"进度回调失败: {e}")
        
        # 如果已提供分析结果，直接使用；否则需要分析
        if editable_images is not None:
            logger.info(f"使用已提供的 {len(editable_images)} 个分析结果创建PPTX")
            report_progress("准备", f"使用已有分析结果（{len(editable_images)} 页）", 10)
        else:
            if not image_paths:
                raise ValueError("必须提供 image_paths 或 editable_images 之一")
            
            total_pages = len(image_paths)
            logger.info(f"开始使用递归分析方法创建可编辑PPTX，共 {total_pages} 页")
            report_progress("开始", f"准备分析 {total_pages} 页幻灯片...", 0)
            
            # 1. 创建ImageEditabilityService（配置自动从 Flask config 获取，使用项目导出设置）
            logger.info(f"使用导出设置: extractor={export_extractor_method}, inpaint={export_inpaint_method}")
            config = ServiceConfig.from_defaults(
                max_depth=max_depth,
                extractor_method=export_extractor_method,
                inpaint_method=export_inpaint_method
            )
            editability_service = ImageEditabilityService(config)
            
            # 2. 并发处理所有页面，生成EditableImage结构
            report_progress("版面分析", f"开始分析 {total_pages} 张图片（并发数: {max_workers}）...", 5)
            from concurrent.futures import ThreadPoolExecutor, as_completed
            
            editable_images = []
            completed_count = 0
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = {
                    executor.submit(editability_service.make_image_editable, img_path): idx
                    for idx, img_path in enumerate(image_paths)
                }
                
                results = [None] * len(image_paths)
                for future in as_completed(futures):
                    idx = futures[future]
                    try:
                        results[idx] = future.result()
                        completed_count += 1
                        # 版面分析占 5% - 40% 的进度
                        percent = 5 + int(35 * completed_count / total_pages)
                        report_progress("版面分析", f"已完成第 {completed_count}/{total_pages} 页的版面分析", percent)
                    except Exception as e:
                        logger.error(f"处理图片 {image_paths[idx]} 失败: {e}")
                        raise
                
                editable_images = results
        
        # 2.5. 使用混合策略提取所有文本元素的样式（如果提供了提取器）
        # 混合策略：全局识别（粗体/斜体/下划线/对齐）+ 单个裁剪识别（颜色）
        text_styles_cache = {}
        if text_attribute_extractor:
            report_progress("样式提取", "开始提取文本样式（混合策略）...", 45)
            
            # 统计文本元素数量
            total_text_count = sum(
                len(ExportService._collect_text_elements_for_extraction(img.elements))
                for img in editable_images
            )
            
            if total_text_count > 0:
                report_progress("样式提取", f"混合策略分析 {total_text_count} 个文本元素...", 50)
                text_styles_cache, failed_extractions = ExportService._batch_extract_text_styles_hybrid(
                    editable_images=editable_images,
                    text_attribute_extractor=text_attribute_extractor,
                    max_workers=max_workers * 2
                )
                
                # 记录样式提取失败的元素（详细）
                for element_id, reason in failed_extractions:
                    warnings.add_style_extraction_failed(element_id, reason)
                
                # 记录汇总信息
                extracted_count = len(text_styles_cache)
                failed_count = len(failed_extractions)
                if failed_count > 0:
                    logger.warning(f"样式提取: {failed_count}/{total_text_count} 个元素失败")
                
                report_progress("样式提取", f"✓ 完成 {extracted_count}/{total_text_count} 个文本样式提取（{failed_count} 个失败）", 70)
        
        report_progress("构建PPTX", "开始构建可编辑PPTX文件...", 75)
        
        # 4. 创建PPTX构建器
        builder = PPTXBuilder()
        builder.create_presentation()
        builder.setup_presentation_size(slide_width_pixels, slide_height_pixels)
        
        # 5. 为每个页面构建幻灯片
        total_pages = len(editable_images)
        for page_idx, editable_img in enumerate(editable_images):
            # 构建PPTX占 75% - 95% 的进度
            percent = 75 + int(20 * page_idx / total_pages)
            report_progress("构建PPTX", f"构建第 {page_idx + 1}/{total_pages} 页...", percent)
            logger.info(f"  构建第 {page_idx + 1}/{total_pages} 页...")
            
            # 创建空白幻灯片
            slide = builder.add_blank_slide()
            
            # 添加背景图（参考原实现，使用slide.shapes.add_picture）
            if editable_img.clean_background and os.path.exists(editable_img.clean_background):
                logger.info(f"    添加clean background: {editable_img.clean_background}")
                try:
                    slide.shapes.add_picture(
                        editable_img.clean_background,
                        left=0,
                        top=0,
                        width=builder.prs.slide_width,
                        height=builder.prs.slide_height
                    )
                except Exception as e:
                    logger.error(f"Failed to add background: {e}")
            else:
                # 回退到原图
                logger.info(f"    使用原图作为背景: {editable_img.image_path}")
                try:
                    slide.shapes.add_picture(
                        editable_img.image_path,
                        left=0,
                        top=0,
                        width=builder.prs.slide_width,
                        height=builder.prs.slide_height
                    )
                except Exception as e:
                    logger.error(f"Failed to add background: {e}")
            
            # 添加所有元素（递归地）
            # 计算缩放比例：将原始图片坐标映射到统一的幻灯片坐标
            # 背景图已经缩放到幻灯片尺寸，所以元素坐标也需要相应缩放
            scale_x = slide_width_pixels / editable_img.width
            scale_y = slide_height_pixels / editable_img.height
            logger.info(f"    元素数量: {len(editable_img.elements)}, 图片尺寸: {editable_img.width}x{editable_img.height}, "
                       f"幻灯片尺寸: {slide_width_pixels}x{slide_height_pixels}, 缩放比例: {scale_x:.3f}x{scale_y:.3f}")
            
            ExportService._add_editable_elements_to_slide(
                builder=builder,
                slide=slide,
                elements=editable_img.elements,
                scale_x=scale_x,
                scale_y=scale_y,
                depth=0,
                text_styles_cache=text_styles_cache,  # 使用预提取的样式缓存
                warnings=warnings  # 收集警告
            )
            
            logger.info(f"    ✓ 第 {page_idx + 1} 页完成，添加了 {len(editable_img.elements)} 个元素")
        
        # 5. 保存或返回字节流
        report_progress("保存文件", "正在保存PPTX文件...", 95)
        if output_file:
            builder.save(output_file)
            report_progress("完成", f"✓ 可编辑PPTX已保存", 100)
            logger.info(f"✓ 可编辑PPTX已保存: {output_file}")
            
            # 输出警告摘要
            if warnings.has_warnings():
                logger.warning(f"导出完成，但有 {len(warnings.to_summary())} 条警告")
            
            return None, warnings
        else:
            pptx_bytes = builder.to_bytes()
            report_progress("完成", f"✓ 可编辑PPTX已生成", 100)
            logger.info(f"✓ 可编辑PPTX已生成（{len(pptx_bytes)} 字节）")
            
            # 输出警告摘要
            if warnings.has_warnings():
                logger.warning(f"导出完成，但有 {len(warnings.to_summary())} 条警告")
            
            return pptx_bytes, warnings
    
    @staticmethod
    def _add_editable_elements_to_slide(
        builder,
        slide,
        elements: List,  # List[EditableElement]
        scale_x: float = 1.0,
        scale_y: float = 1.0,
        depth: int = 0,
        text_styles_cache: Dict[str, Any] = None,  # 预提取的文本样式缓存，key为element_id
        warnings: 'ExportWarnings' = None  # 警告收集器
    ):
        """
        递归地将EditableElement添加到幻灯片
        
        Args:
            builder: PPTXBuilder实例
            slide: 幻灯片对象
            elements: EditableElement列表
            scale_x: X轴缩放因子
            scale_y: Y轴缩放因子
            depth: 当前递归深度
            text_styles_cache: 预提取的文本样式缓存（可选），由 _batch_extract_text_styles 生成
        
        Note:
            elem.image_path 现在是绝对路径，无需额外的目录参数
        """
        if text_styles_cache is None:
            text_styles_cache = {}
        
        for elem in elements:
            elem_type = elem.element_type
            
            # 根据深度决定使用局部坐标还是全局坐标
            # depth=0: 顶层元素，使用局部坐标（bbox）
            # depth>0: 子元素，需要使用全局坐标（bbox_global）
            if depth == 0:
                bbox = elem.bbox  # 顶层元素使用局部坐标
            else:
                bbox = elem.bbox_global if hasattr(elem, 'bbox_global') and elem.bbox_global else elem.bbox
            
            # 转换BBox对象为列表并应用缩放
            bbox_list = [
                int(bbox.x0 * scale_x),
                int(bbox.y0 * scale_y),
                int(bbox.x1 * scale_x),
                int(bbox.y1 * scale_y)
            ]
            
            logger.info(f"{'  ' * depth}  添加元素: type={elem_type}, bbox={bbox_list}, content={elem.content[:30] if elem.content else None}, image_path={elem.image_path}, 使用{'全局' if depth > 0 else '局部'}坐标")
            
            # 根据类型添加元素（参考原实现的_add_mineru_text_to_slide和_add_mineru_image_to_slide）
            if elem_type in ['text', 'title', 'list', 'paragraph', 'header', 'footer', 'heading', 'table_caption', 'image_caption']:
                # 添加文本（参考_add_mineru_text_to_slide）
                if elem.content:
                    text = elem.content.strip()
                    if text:
                        try:
                            # 确定文本级别
                            level = 'title' if elem_type in ['title', 'heading'] else 'default'
                            
                            # 从缓存获取预提取的文字样式
                            text_style = text_styles_cache.get(elem.element_id)
                            if text_style:
                                logger.debug(f"{'  ' * depth}  使用缓存的文字样式: color={text_style.font_color_rgb}, bold={text_style.is_bold}")
                            
                            builder.add_text_element(
                                slide=slide,
                                text=text,
                                bbox=bbox_list,
                                text_level=level,
                                text_style=text_style
                            )
                        except Exception as e:
                            logger.warning(f"添加文本元素失败: {e}")
                            if warnings:
                                warnings.add_text_render_failed(text, str(e))
            
            elif elem_type == 'table_cell':
                # 添加表格单元格（带边框的文本框）
                if elem.content:
                    text = elem.content.strip()
                    if text:
                        try:
                            # 从缓存获取预提取的文字样式
                            text_style = text_styles_cache.get(elem.element_id)
                            
                            # 表格单元格已经在上面统一处理了bbox_global和缩放
                            # 直接使用bbox_list即可
                            builder.add_text_element(
                                slide=slide,
                                text=text,
                                bbox=bbox_list,
                                text_level=None,
                                align='center',
                                text_style=text_style
                            )
                            
                        except Exception as e:
                            logger.warning(f"添加单元格失败: {e}")
                            if warnings:
                                warnings.add_text_render_failed(text, str(e))
            
            elif elem_type == 'table':
                # 如果表格有子元素（单元格），使用inpainted背景 + 单元格
                if elem.children and elem.inpainted_background_path:
                    logger.info(f"{'  ' * depth}    表格有 {len(elem.children)} 个单元格，使用可编辑格式")
                    
                    # 先添加inpainted背景（干净的表格框架）
                    if os.path.exists(elem.inpainted_background_path):
                        try:
                            builder.add_image_element(
                                slide=slide,
                                image_path=elem.inpainted_background_path,
                                bbox=bbox_list
                            )
                        except Exception as e:
                            logger.error(f"Failed to add table background: {e}")
                    
                    # 递归添加单元格
                    ExportService._add_editable_elements_to_slide(
                        builder=builder,
                        slide=slide,
                        elements=elem.children,
                        scale_x=scale_x,
                        scale_y=scale_y,
                        depth=depth + 1,
                        text_styles_cache=text_styles_cache,
                        warnings=warnings
                    )
                else:
                    # 没有子元素，添加整体表格图片
                    # elem.image_path 现在是绝对路径
                    if elem.image_path and os.path.exists(elem.image_path):
                        try:
                            builder.add_image_element(
                                slide=slide,
                                image_path=elem.image_path,
                                bbox=bbox_list
                            )
                        except Exception as e:
                            logger.error(f"Failed to add table image: {e}")
                    else:
                        logger.warning(f"Table image not found: {elem.image_path}")
                        builder.add_image_placeholder(slide, bbox_list)
            
            elif elem_type in ['image', 'figure', 'chart']:
                # 检查是否应该使用递归渲染
                should_use_recursive_render = False
                
                if elem.children and elem.inpainted_background_path:
                    # 检查是否有任意子元素占据父元素绝大部分面积
                    parent_area = (bbox.x1 - bbox.x0) * (bbox.y1 - bbox.y0)
                    max_child_coverage_ratio = 0.85  # 阈值
                    has_dominant_child = False
                    
                    for child in elem.children:
                        if hasattr(child, 'bbox_global') and child.bbox_global:
                            child_bbox = child.bbox_global
                        else:
                            child_bbox = child.bbox
                        
                        child_area = child_bbox.area
                        coverage_ratio = child_area / parent_area if parent_area > 0 else 0
                        
                        if coverage_ratio > max_child_coverage_ratio:
                            logger.info(f"{'  ' * depth}    子元素 {child.element_id} 占父元素面积 {coverage_ratio*100:.1f}% (>{max_child_coverage_ratio*100:.0f}%)，跳过递归渲染，直接使用原图")
                            has_dominant_child = True
                            break
                    
                    should_use_recursive_render = not has_dominant_child
                
                # 如果有子元素且应该递归渲染
                if should_use_recursive_render:
                    logger.debug(f"{'  ' * depth}    元素有 {len(elem.children)} 个子元素，递归添加")
                    
                    # 先添加inpainted背景
                    if os.path.exists(elem.inpainted_background_path):
                        try:
                            builder.add_image_element(slide, elem.inpainted_background_path, bbox_list)
                        except Exception as e:
                            logger.error(f"Failed to add inpainted background: {e}")
                    
                    # 递归添加子元素
                    ExportService._add_editable_elements_to_slide(
                        builder=builder,
                        slide=slide,
                        elements=elem.children,
                        scale_x=scale_x,
                        scale_y=scale_y,
                        depth=depth + 1,
                        text_styles_cache=text_styles_cache,
                        warnings=warnings
                    )
                else:
                    # 没有子元素或子元素占比过大，直接添加原图
                    # elem.image_path 现在是绝对路径
                    if elem.image_path and os.path.exists(elem.image_path):
                        try:
                            builder.add_image_element(
                                slide=slide,
                                image_path=elem.image_path,
                                bbox=bbox_list
                            )
                        except Exception as e:
                            logger.error(f"Failed to add image: {e}")
                    else:
                        logger.warning(f"Image file not found: {elem.image_path}")
                        builder.add_image_placeholder(slide, bbox_list)
            
            else:
                # 其他类型
                logger.debug(f"{'  ' * depth}  跳过未知类型: {elem_type}")
    

